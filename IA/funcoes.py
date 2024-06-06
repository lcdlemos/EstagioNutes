from pptx.util import Inches
import os
import requests
import json
from idkey import id_key
import urllib.request
from PIL import Image
from senha import API_KEY

def contar_slides(nome_arquivo):
    with open(nome_arquivo, "r", encoding="utf-8") as arquivo:
        texto = arquivo.read()
        num_slides = texto.count("Slide")
    return num_slides

def ler_arquivo(nome_arquivo):
    with open(nome_arquivo, "r", encoding="utf-8") as arquivo:
        texto = arquivo.readlines()
    return texto

def criar_capa(auto_apresentacao, texto):
    modelo_slide = auto_apresentacao.slide_layouts[0]
    slide = auto_apresentacao.slides.add_slide(modelo_slide)
    titulo = slide.shapes.title
    subtitulo = slide.placeholders[1]
    if "Título" in texto[1]:
        titulo.text = texto[1].lstrip("- Título: ").rstrip("\n")
        subtitulo.text = texto[2].strip("\n")
    elif "Tema" in texto[1]:
        titulo.text = texto[1].lstrip("Tema: ").rstrip("\n")
        subtitulo.text = texto[2].strip("\n")
    else:
        titulo.text = texto[0].lstrip("Slide 0: ").rstrip("\n")
        subtitulo.text = texto[1].strip("\n")
    return slide

def coletar_titulos(texto):
    titulos = []
    for linha in texto:
        if "Título" in linha:
            titulos.append(linha.lstrip("- Título: ").rstrip("\n"))
    if titulos == []:
        for linha in texto:
            if "Slide" in linha:
                titulos.append(linha.lstrip("Slide ").rstrip("\n"))
    return titulos

def coletar_topicos(texto):
    topicos = []
    for topico in texto:
        if "- " in topico and "Título" not in topico and "Autor:" not in topico or "\n" == topico:
            topicos.append(topico)
    topicos.append("\n")
    return topicos

def adicionar_slides(auto_apresentacao, texto, titulos, topicos, num_slides):
    slides = []
    j = 1
    for i in range(num_slides - 1):
        modelo_slide = auto_apresentacao.slide_layouts[1]
        slide = auto_apresentacao.slides.add_slide(modelo_slide)
        titulo = slide.shapes.title
        titulo.text = titulos[i+1].lstrip("- Título: ").rstrip("\n")
        caixa_topicos = slide.shapes
        caixa_topico = caixa_topicos.placeholders[1]
        caixa_topico.text = topicos[j].strip("- ""\n")
        j += 1
        while topicos[j] != "\n":
            if "    -" in topicos[j]:
                topico = caixa_topico.text_frame.add_paragraph()
                topico.text = topicos[j].strip("    -""\n")
                topico.level = 1
            elif "Imagem" in topicos[j]:
                imagem = r"C:\Users\Luiz Carlos\Documents\UEPB\Estágio\Slides\ia.jpg"
                if os.path.isfile(imagem):
                    img = "ia.jpg"
                    esquerda = Inches(3)
                    topo = Inches(4.5)
                    adiciona_imagem = slide.shapes.add_picture(img, esquerda, topo)
                else:
                    #BUSCA NO UNSPLASH
                    #busca = "beach"
                    #requisicao = requests.get(f"https://api.unsplash.com/search/photos?client_id={id_key}&query={busca}")
                    
                    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}
                    id_modelo = "dall-e-3"
                    link = "https://api.openai.com/v1/images/generations"

                    body_message = {
                        "model": id_modelo,
                        "prompt": titulos[i],
                        "n": 1,
                        "size": "1024x1024",
                    }

                    body_message = json.dumps(body_message)
                    requisicao = requests.post(link, headers=headers, data=body_message)

                    resposta = requisicao.json()
                    #url_image = resposta['results'][0]['urls']['small']
                    url_image = resposta['data'][0]['url']
                    urllib.request.urlretrieve(url_image, "./imagem_url.jpeg")
                    
                    ima = Image.open("imagem_url.jpeg")
                    ima = ima.resize((300, 200))
                    ima.save("imagem_url.jpeg")
                    img = "imagem_url.jpeg"
                    esquerda = Inches(3)
                    topo = Inches(4.5)
                    adiciona_imagem = slide.shapes.add_picture(img, esquerda, topo)
            else:
                topico = caixa_topico.text_frame.add_paragraph()
                topico.text = topicos[j].strip("- ""\n")
            j += 1
        j += 1
        slides.append(slide)
    return slides